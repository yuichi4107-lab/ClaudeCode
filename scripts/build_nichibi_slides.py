"""日美株式会社の会社紹介スライド (.pptx) を生成するスクリプト。

実行:
    python scripts/build_nichibi_slides.py

出力: docs/slides/nichibi_company_intro.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


PRIMARY = RGBColor(0x0B, 0x3D, 0x91)
ACCENT = RGBColor(0xE8, 0x1C, 0x4F)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT = RGBColor(0x1B, 0x1F, 0x2A)
SUB_TEXT = RGBColor(0x4A, 0x55, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER = RGBColor(0xD7, 0xDC, 0xE3)

JP_FONT = "Yu Gothic"
JP_FONT_BOLD = "Yu Gothic"

SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_fill(shape, color)
    return shape


def _add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = DARK_TEXT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font: str = JP_FONT,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0)
    tf.margin_right = Cm(0)
    tf.margin_top = Cm(0)
    tf.margin_bottom = Cm(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def _add_header(slide, title: str, subtitle: str | None = None) -> None:
    # サイドバー
    _add_rect(slide, Cm(0), Cm(0), Cm(0.5), SLIDE_H, PRIMARY)
    # 見出し
    _add_text(
        slide,
        Cm(1.2),
        Cm(0.8),
        Cm(30),
        Cm(1.5),
        title,
        size=28,
        bold=True,
        color=PRIMARY,
    )
    if subtitle:
        _add_text(
            slide,
            Cm(1.2),
            Cm(2.2),
            Cm(30),
            Cm(0.8),
            subtitle,
            size=14,
            color=SUB_TEXT,
        )
    # 区切り線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(1.2), Cm(3.2), Cm(31), Cm(0.05)
    )
    _set_fill(line, DIVIDER)


def _add_footer(slide, page_no: int, total: int) -> None:
    _add_text(
        slide,
        Cm(1.2),
        Cm(18.1),
        Cm(20),
        Cm(0.6),
        "日美株式会社  Company Profile",
        size=9,
        color=SUB_TEXT,
    )
    _add_text(
        slide,
        Cm(30),
        Cm(18.1),
        Cm(3),
        Cm(0.6),
        f"{page_no} / {total}",
        size=9,
        color=SUB_TEXT,
        align=PP_ALIGN.RIGHT,
    )


def _add_bullets(slide, left, top, width, height, items: list[str], *, size: int = 16):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0)
    tf.margin_right = Cm(0)
    tf.margin_top = Cm(0)
    tf.margin_bottom = Cm(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"●  {item}"
        run.font.name = JP_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = DARK_TEXT


# ---------- 各スライド ----------


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景
    _add_rect(slide, Cm(0), Cm(0), SLIDE_W, SLIDE_H, PRIMARY)
    # アクセント帯
    _add_rect(slide, Cm(0), Cm(13.5), SLIDE_W, Cm(0.15), ACCENT)
    # ロゴ風円
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(2.5), Cm(5.5), Cm(2.4), Cm(2.4))
    _set_fill(circ, WHITE)
    _add_text(
        slide,
        Cm(2.5),
        Cm(5.5),
        Cm(2.4),
        Cm(2.4),
        "N",
        size=44,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # タイトル
    _add_text(
        slide,
        Cm(5.5),
        Cm(5.6),
        Cm(25),
        Cm(2),
        "日美株式会社",
        size=54,
        bold=True,
        color=WHITE,
    )
    _add_text(
        slide,
        Cm(5.5),
        Cm(7.8),
        Cm(25),
        Cm(1),
        "Nichibi Co., Ltd.   Company Profile",
        size=20,
        color=RGBColor(0xCD, 0xD8, 0xF0),
    )
    # サブコピー
    _add_text(
        slide,
        Cm(2.5),
        Cm(11.0),
        Cm(28),
        Cm(2),
        "人と企業の未来をつなぐ。\n名古屋発、東海エリアの人材パートナー。",
        size=22,
        color=WHITE,
    )
    _add_text(
        slide,
        Cm(2.5),
        Cm(17.5),
        Cm(28),
        Cm(1),
        "https://nichibi-web.co.jp",
        size=14,
        color=RGBColor(0xCD, 0xD8, 0xF0),
    )


def slide_agenda(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "目次", "Agenda")
    items = [
        "1.  会社概要",
        "2.  私たちのミッション",
        "3.  事業内容",
        "4.  取扱職種",
        "5.  対応エリア",
        "6.  日美の強み",
        "7.  サポート体制",
        "8.  ご利用の流れ",
        "9.  お問い合わせ",
    ]
    # 2カラム
    left_items = items[:5]
    right_items = items[5:]
    for i, t in enumerate(left_items):
        _add_text(
            slide,
            Cm(2.5),
            Cm(4.5 + i * 1.6),
            Cm(14),
            Cm(1.2),
            t,
            size=20,
            color=DARK_TEXT,
        )
    for i, t in enumerate(right_items):
        _add_text(
            slide,
            Cm(18),
            Cm(4.5 + i * 1.6),
            Cm(14),
            Cm(1.2),
            t,
            size=20,
            color=DARK_TEXT,
        )
    _add_footer(slide, page, total)


def slide_company(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "会社概要", "Company Overview")

    rows = [
        ("社名", "日美株式会社 (Nichibi Co., Ltd.)"),
        ("設立", "1984年5月"),
        ("代表者", "代表取締役  長谷川 裕一"),
        ("本社所在地", "愛知県名古屋市中区錦1丁目5番13号"),
        ("資本金", "50,000,000 円"),
        ("従業員数", "25名"),
        ("事業内容", "人材派遣 / アウトソーシング / 人材紹介"),
        ("Web サイト", "https://nichibi-web.co.jp"),
    ]
    top = Cm(4.2)
    row_h = Cm(1.3)
    label_w = Cm(6.5)
    value_w = Cm(22)

    for i, (label, value) in enumerate(rows):
        y = top + row_h * i
        # ラベル背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Cm(2.5), y, label_w, row_h
        )
        _set_fill(bg, LIGHT_BG)
        _add_text(
            slide,
            Cm(2.5),
            y,
            label_w,
            row_h,
            "  " + label,
            size=14,
            bold=True,
            color=PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_text(
            slide,
            Cm(2.5) + label_w + Cm(0.4),
            y,
            value_w,
            row_h,
            value,
            size=14,
            color=DARK_TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    _add_footer(slide, page, total)


def slide_mission(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "私たちのミッション", "Our Mission")

    _add_text(
        slide,
        Cm(2.5),
        Cm(4.5),
        Cm(29),
        Cm(2.5),
        "人と企業の架け橋として、\n働く人の「ありたい姿」と企業の「成長」を共に実現する。",
        size=26,
        bold=True,
        color=PRIMARY,
    )

    boxes = [
        ("一人ひとりに寄り添う", "声を聴き、ご希望に沿う働き方をとことんサポート。"),
        ("企業の成長を支える", "ニーズにフレキシブルに応える人材ソリューション。"),
        ("地域に根差す", "愛知・岐阜・三重・静岡で地域密着のネットワーク。"),
    ]
    box_w = Cm(9.7)
    gap = Cm(0.6)
    box_h = Cm(6.5)
    top = Cm(9.5)
    for i, (h, body) in enumerate(boxes):
        left = Cm(2.5) + (box_w + gap) * i
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        _set_fill(bg, LIGHT_BG)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_w, Cm(0.25))
        _set_fill(bar, ACCENT)
        _add_text(
            slide,
            left + Cm(0.4),
            top + Cm(0.6),
            box_w - Cm(0.8),
            Cm(1.4),
            h,
            size=18,
            bold=True,
            color=PRIMARY,
        )
        _add_text(
            slide,
            left + Cm(0.4),
            top + Cm(2.4),
            box_w - Cm(0.8),
            Cm(4),
            body,
            size=13,
            color=DARK_TEXT,
        )
    _add_footer(slide, page, total)


def slide_business(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "事業内容", "Our Business")

    items = [
        ("人材派遣", "正社員の代替要員・欠員補充、繁忙期対応など。\n貴社のニーズに応じて即戦力人材をご紹介します。"),
        ("アウトソーシング", "業務プロセスごとの受託で、コア業務への\n集中を支援。生産性向上と固定費の最適化に貢献。"),
        ("人材紹介", "登録者2万人超のネットワークから、\n貴社にマッチする人材をスピーディーにご提案。"),
    ]
    box_w = Cm(9.7)
    gap = Cm(0.6)
    box_h = Cm(11)
    top = Cm(4.5)
    for i, (h, body) in enumerate(items):
        left = Cm(2.5) + (box_w + gap) * i
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        _set_fill(bg, WHITE)
        bg.line.color.rgb = DIVIDER
        # ヘッダ帯
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_w, Cm(2))
        _set_fill(head, PRIMARY)
        _add_text(
            slide,
            left,
            top,
            box_w,
            Cm(2),
            f"0{i+1}",
            size=16,
            bold=True,
            color=RGBColor(0xCD, 0xD8, 0xF0),
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.TOP,
        )
        _add_text(
            slide,
            left,
            top + Cm(0.6),
            box_w,
            Cm(1.4),
            h,
            size=20,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_text(
            slide,
            left + Cm(0.6),
            top + Cm(2.6),
            box_w - Cm(1.2),
            box_h - Cm(3),
            body,
            size=14,
            color=DARK_TEXT,
        )
    _add_footer(slide, page, total)


def slide_jobs(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "取扱職種", "Job Categories")

    categories = [
        (
            "事務系",
            ["一般事務", "営業事務", "経理・財務", "人事・総務", "貿易事務", "金融事務（銀行/証券/損保/生保）", "受発注事務", "受付・案内"],
        ),
        (
            "専門・技術系",
            ["CADオペレーター", "インテリアコーディネーター", "秘書", "テレホンオペレーター", "保険関連業務"],
        ),
        (
            "医療・福祉系",
            ["医療事務", "医師事務作業補助", "看護補助", "介護", "各種検査技師"],
        ),
        (
            "IT・エンジニア系",
            ["システムエンジニア", "プログラマー", "ネットワークエンジニア", "ヘルプデスク", "ユーザー / プロバイダーサポート", "アドミニストレーター"],
        ),
    ]

    box_w = Cm(14.8)
    box_h = Cm(6.4)
    gap_x = Cm(0.6)
    gap_y = Cm(0.5)
    origin_l = Cm(2.5)
    origin_t = Cm(4.4)

    for idx, (cat, jobs) in enumerate(categories):
        r, c = divmod(idx, 2)
        left = origin_l + (box_w + gap_x) * c
        top = origin_t + (box_h + gap_y) * r
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        _set_fill(bg, LIGHT_BG)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Cm(0.25), box_h)
        _set_fill(bar, PRIMARY)
        _add_text(
            slide,
            left + Cm(0.6),
            top + Cm(0.4),
            box_w - Cm(1),
            Cm(1),
            cat,
            size=16,
            bold=True,
            color=PRIMARY,
        )
        body = "  /  ".join(jobs)
        _add_text(
            slide,
            left + Cm(0.6),
            top + Cm(1.6),
            box_w - Cm(1),
            box_h - Cm(2),
            body,
            size=12,
            color=DARK_TEXT,
        )
    _add_footer(slide, page, total)


def slide_area(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "対応エリア", "Service Area")

    _add_text(
        slide,
        Cm(2.5),
        Cm(4.3),
        Cm(29),
        Cm(1.2),
        "東海4県を中心に、地域に根差したサポートを提供します。",
        size=18,
        color=DARK_TEXT,
    )

    areas = [
        ("愛知", "本社・名古屋拠点"),
        ("岐阜", "東海エリア対応"),
        ("三重", "東海エリア対応"),
        ("静岡", "東海エリア対応"),
    ]
    box_w = Cm(7.0)
    gap = Cm(0.6)
    box_h = Cm(7.5)
    top = Cm(6.8)
    total_w = box_w * 4 + gap * 3
    start_l = (SLIDE_W - total_w) / 2
    for i, (name, sub) in enumerate(areas):
        left = start_l + (box_w + gap) * i
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        color = PRIMARY if i == 0 else LIGHT_BG
        _set_fill(bg, color)
        _add_text(
            slide,
            left,
            top + Cm(1.5),
            box_w,
            Cm(3),
            name,
            size=44,
            bold=True,
            color=WHITE if i == 0 else PRIMARY,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_text(
            slide,
            left,
            top + Cm(5.0),
            box_w,
            Cm(1.5),
            sub,
            size=12,
            color=WHITE if i == 0 else SUB_TEXT,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    _add_footer(slide, page, total)


def slide_strength(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "日美の強み", "Our Strengths")

    strengths = [
        ("01", "40年以上の実績", "1984年創業。\n東海エリアでの長年の信頼と実績。"),
        ("02", "登録者2万人超", "創業以来の登録者は2万人以上。\n有資格者・経験者も多数。"),
        ("03", "幅広い職種対応", "事務・専門職・医療・ITまで\n多彩な職種をカバー。"),
        ("04", "地域密着ネットワーク", "愛知・岐阜・三重・静岡で\nきめ細かなマッチングを実現。"),
        ("05", "手厚いサポート", "派遣先への定期訪問・\nストレスチェックなど就業中も安心。"),
        ("06", "柔軟な提案力", "派遣・紹介・アウトソーシングを\n組み合わせ最適解をご提案。"),
    ]
    box_w = Cm(9.7)
    box_h = Cm(6.5)
    gap_x = Cm(0.6)
    gap_y = Cm(0.5)
    origin_l = Cm(2.5)
    origin_t = Cm(4.3)
    for i, (num, h, body) in enumerate(strengths):
        r, c = divmod(i, 3)
        left = origin_l + (box_w + gap_x) * c
        top = origin_t + (box_h + gap_y) * r
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        _set_fill(bg, WHITE)
        bg.line.color.rgb = DIVIDER
        _add_text(
            slide,
            left + Cm(0.6),
            top + Cm(0.4),
            Cm(2),
            Cm(1.2),
            num,
            size=28,
            bold=True,
            color=ACCENT,
        )
        _add_text(
            slide,
            left + Cm(0.6),
            top + Cm(2.0),
            box_w - Cm(1.2),
            Cm(1.3),
            h,
            size=16,
            bold=True,
            color=PRIMARY,
        )
        _add_text(
            slide,
            left + Cm(0.6),
            top + Cm(3.5),
            box_w - Cm(1.2),
            Cm(3),
            body,
            size=12,
            color=DARK_TEXT,
        )
    _add_footer(slide, page, total)


def slide_support(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "サポート体制", "Our Support")

    _add_text(
        slide,
        Cm(2.5),
        Cm(4.3),
        Cm(29),
        Cm(1.2),
        "登録から就業後まで、専任スタッフが伴走します。",
        size=18,
        color=DARK_TEXT,
    )

    items = [
        ("登録前カウンセリング", "希望条件・キャリアプランを丁寧にヒアリング。"),
        ("マッチングのご提案", "登録者のスキルと企業ニーズを精緻にマッチング。"),
        ("就業前オリエンテーション", "業務内容・職場環境を事前にご説明し、不安を解消。"),
        ("派遣先への定期訪問", "就業状況を定期確認し、お困りごとに迅速に対応。"),
        ("ストレスチェック", "心身の健康を継続フォロー。長く安心して働ける環境を。"),
        ("キャリア相談", "次のステップ・スキルアップのご相談にも対応。"),
    ]
    col_w = Cm(14.8)
    row_h = Cm(1.7)
    gap_x = Cm(0.6)
    top = Cm(6.5)
    for i, (h, body) in enumerate(items):
        r, c = divmod(i, 2)
        left = Cm(2.5) + (col_w + gap_x) * c
        y = top + row_h * r
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y, Cm(1.3), Cm(1.3))
        _set_fill(circle, PRIMARY)
        _add_text(
            slide,
            left,
            y,
            Cm(1.3),
            Cm(1.3),
            f"{i+1:02d}",
            size=12,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_text(
            slide,
            left + Cm(1.6),
            y,
            col_w - Cm(1.6),
            Cm(0.8),
            h,
            size=14,
            bold=True,
            color=PRIMARY,
        )
        _add_text(
            slide,
            left + Cm(1.6),
            y + Cm(0.75),
            col_w - Cm(1.6),
            Cm(0.9),
            body,
            size=11,
            color=DARK_TEXT,
        )
    _add_footer(slide, page, total)


def slide_flow(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "ご利用の流れ", "How it works")

    steps = [
        ("STEP 1", "お問い合わせ", "Webまたはお電話でご相談ください。"),
        ("STEP 2", "ヒアリング", "ご要望や条件を詳しくお伺いします。"),
        ("STEP 3", "ご提案", "最適な人材・プランをご提案します。"),
        ("STEP 4", "就業開始", "オリエンテーション後、業務スタート。"),
        ("STEP 5", "アフターフォロー", "定期訪問・面談で継続的にサポート。"),
    ]
    box_w = Cm(5.7)
    box_h = Cm(6.0)
    gap = Cm(0.45)
    top = Cm(6.5)
    total_w = box_w * 5 + gap * 4
    start_l = (SLIDE_W - total_w) / 2
    for i, (step, title, body) in enumerate(steps):
        left = start_l + (box_w + gap) * i
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        _set_fill(bg, LIGHT_BG)
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_w, Cm(1.2))
        _set_fill(head, PRIMARY)
        _add_text(
            slide,
            left,
            top,
            box_w,
            Cm(1.2),
            step,
            size=12,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_text(
            slide,
            left + Cm(0.3),
            top + Cm(1.5),
            box_w - Cm(0.6),
            Cm(1.2),
            title,
            size=16,
            bold=True,
            color=PRIMARY,
            align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide,
            left + Cm(0.3),
            top + Cm(3.0),
            box_w - Cm(0.6),
            Cm(3),
            body,
            size=11,
            color=DARK_TEXT,
            align=PP_ALIGN.CENTER,
        )
        # 矢印
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                left + box_w + Cm(0.02),
                top + box_h / 2 - Cm(0.3),
                Cm(0.4),
                Cm(0.6),
            )
            _set_fill(arrow, ACCENT)
    _add_footer(slide, page, total)


def slide_contact(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Cm(0), Cm(0), SLIDE_W, SLIDE_H, PRIMARY)
    _add_rect(slide, Cm(0), Cm(4.0), SLIDE_W, Cm(0.1), ACCENT)

    _add_text(
        slide,
        Cm(2.5),
        Cm(1.4),
        Cm(28),
        Cm(2),
        "お問い合わせ",
        size=40,
        bold=True,
        color=WHITE,
    )
    _add_text(
        slide,
        Cm(2.5),
        Cm(3.0),
        Cm(28),
        Cm(1),
        "Contact Us",
        size=16,
        color=RGBColor(0xCD, 0xD8, 0xF0),
    )

    info = [
        ("会社名", "日美株式会社"),
        ("本社", "愛知県名古屋市中区錦1丁目5番13号"),
        ("Web", "https://nichibi-web.co.jp"),
        ("お問い合わせ", "Webサイトのお問い合わせフォームよりご連絡ください"),
        ("対応エリア", "愛知 / 岐阜 / 三重 / 静岡"),
    ]
    top = Cm(6.0)
    for i, (k, v) in enumerate(info):
        y = top + Cm(1.4) * i
        _add_text(
            slide,
            Cm(2.5),
            y,
            Cm(7),
            Cm(1.2),
            k,
            size=14,
            bold=True,
            color=RGBColor(0xCD, 0xD8, 0xF0),
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_text(
            slide,
            Cm(9.5),
            y,
            Cm(22),
            Cm(1.2),
            v,
            size=18,
            color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    _add_text(
        slide,
        Cm(2.5),
        Cm(16.5),
        Cm(28),
        Cm(1.5),
        "人と企業の未来をつなぐパートナーとして、お気軽にご相談ください。",
        size=16,
        color=RGBColor(0xCD, 0xD8, 0xF0),
    )
    _add_text(
        slide,
        Cm(2.5),
        Cm(18.0),
        Cm(28),
        Cm(0.8),
        f"{page} / {total}",
        size=9,
        color=RGBColor(0xCD, 0xD8, 0xF0),
        align=PP_ALIGN.RIGHT,
    )


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_agenda,
        slide_company,
        slide_mission,
        slide_business,
        slide_jobs,
        slide_area,
        slide_strength,
        slide_support,
        slide_flow,
        slide_contact,
    ]
    total = len(builders)

    for i, build in enumerate(builders, start=1):
        if build is slide_title:
            build(prs)
        else:
            build(prs, i, total)

    out_dir = Path(__file__).resolve().parent.parent / "docs" / "slides"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "nichibi_company_intro.pptx"
    prs.save(out_path)
    print(f"saved: {out_path}")


if __name__ == "__main__":
    main()
